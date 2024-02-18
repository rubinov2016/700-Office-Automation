import os
import sys
from pptx import Presentation
# pip install python-pptx


def extract_text_from_pptx_file(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def extract_text_from_folder(folder_path):
    extracted_text = []
    for root, dirs, files in os.walk(folder_path):
        for file_name in files:
            if file_name.endswith('.pptx'):
                file_path = os.path.join(root, file_name)
                text = extract_text_from_pptx_file(file_path)
                extracted_text.append(text)
    return '\n'.join(extracted_text)

# Replace 'your_folder_path' with the path to your folder containing PowerPoint files
# folder_path = r'C:\Users\Lenovo\OneDrive - Solent University\Documents\724 AI in Business (Friday)\Lesson 1\'
folder_path = r'C:\Users\Lenovo\OneDrive - Solent University\Documents\724 AI in Business (Friday)\Lesson 1'
all_texts = extract_text_from_folder(folder_path)
print(all_texts.encode('utf-8').decode(sys.stdout.encoding, 'ignore'))

